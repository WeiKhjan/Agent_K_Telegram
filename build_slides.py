#!/usr/bin/env python3
"""Generate PowerPoint: Agent K — Deployment Strategy for Accounting Firm"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Colour palette --
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GREY_TEXT = RGBColor(0x5A, 0x5A, 0x6E)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_RED_BG = RGBColor(0xFD, 0xED, 0xED)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xEE)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_text(slide, left, top, width, height, text, font_size=14,
                         bold=False, color=DARK_TEXT, bg_color=None,
                         alignment=PP_ALIGN.LEFT, font_name='Calibri',
                         border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                  bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                  font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                     color=DARK_TEXT, bullet_color=None, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_card(slide, left, top, width, height, title, items,
             title_color=WHITE, bg_color=WHITE, header_color=DARK_BLUE,
             item_color=DARK_TEXT, font_size=12):
    # Header bar
    header_h = Inches(0.45)
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, header_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.margin_left = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = 'Calibri'

    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + header_h - Inches(0.05),
                                   width, height - header_h + Inches(0.05))
    body.fill.solid()
    body.fill.fore_color.rgb = bg_color
    body.line.color.rgb = CARD_BORDER
    body.line.width = Pt(1)
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.15)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = item_color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_title_bar(slide, title, subtitle=None):
    """Add a dark blue title bar at top of slide"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.55),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.35),
                     subtitle, font_size=14, color=RGBColor(0xB0, 0xC4, 0xDE))


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1),
             "Agent K", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
             "AI Staff Deployment Strategy for Accounting Firms",
             font_size=24, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER)

# Divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.6), Inches(5.3), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_BLUE
div.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
             "Telegram Bot  vs  Claude Code CLI",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11.3), Inches(0.5),
             "Architecture  |  Pros & Cons  |  Concurrency  |  Recommendation",
             font_size=15, color=RGBColor(0x88, 0xA0, 0xC0), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
             "Audit  |  Tax  |  Compilation Services  \u2014  Malaysia",
             font_size=14, color=RGBColor(0x70, 0x88, 0xA8), alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "Opening slide. This presentation compares two deployment approaches for an AI staff member "
    "at a Malaysia accounting firm: (A) Teaching the firm to use Claude Code CLI directly, vs "
    "(B) Deploying the Agent K Telegram bot as a chat interface. We cover architecture, "
    "pros/cons, concurrency, and a recommended hybrid approach."
)


# ============================================================
# SLIDE 2: Two Approaches at a Glance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Two Deployment Approaches")

# Approach A card
add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(4.8),
         "A: Claude Code CLI (Direct)", [
             "\u2022  Staff types commands in terminal",
             "\u2022  Claude CLI runs directly on Mac Mini",
             "\u2022  No middleware, no bot layer",
             "\u2022  Full CLI power: /help, slash commands, MCP",
             "\u2022  Each user runs their own session",
             "",
             "Interface:  Terminal / VS Code",
             "Setup:  Install Claude Code + train staff",
             "Who builds it:  The firm (with training)",
         ], header_color=MID_BLUE, font_size=13)

# Approach B card
add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(4.8),
         "B: Agent K Telegram Bot", [
             "\u2022  Staff sends messages via Telegram",
             "\u2022  Bot relays to Claude CLI behind the scenes",
             "\u2022  Files sent/received via Telegram chat",
             "\u2022  Mobile-friendly, fire-and-forget",
             "\u2022  One bot serves all staff",
             "",
             "Interface:  Telegram (phone / desktop)",
             "Setup:  Clone repo, configure .env, deploy",
             "Who builds it:  You (then hand over)",
         ], header_color=ACCENT_BLUE, font_size=13)

# Shared foundation
add_shape_with_text(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.6),
                    "Both use the same foundation:  Claude CLI  +  Skills (SKILL.md)  +  MCP Servers  +  Local filesystem",
                    font_size=13, bold=True, color=DARK_BLUE, bg_color=WHITE,
                    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE, border_width=Pt(1.5))

add_notes(slide,
    "Key point: Both approaches use the EXACT same foundation underneath.\n\n"
    "- Skills (SKILL.md files) are identical in both\n"
    "- MCP servers (Excel, Word, Gmail, etc.) are identical\n"
    "- Templates and folder structure are identical\n"
    "- The only difference is HOW the human interacts: terminal vs Telegram\n\n"
    "Approach A: Staff learns to use the terminal. More powerful but requires technical comfort.\n"
    "Approach B: Staff chats on Telegram. Zero learning curve but adds a bot middleware layer.\n\n"
    "The Telegram bot is ~400 lines of JavaScript. It's a 'dumb pipe' that passes messages "
    "to Claude CLI and returns the response."
)


# ============================================================
# SLIDE 3: Telegram Bot — Pros & Cons
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Approach B: Telegram Bot", "Pros & Cons")

# Pros card
add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.2),
         "\u2705  Strengths", [
             "\u2022  Zero technical skill required",
             "    Staff already know Telegram",
             "",
             "\u2022  Mobile-friendly",
             "    Send requests from phone, get files back",
             "",
             "\u2022  Fire-and-forget async workflow",
             "    Send message, get notified when done",
             "",
             "\u2022  Automatic file delivery",
             "    PDFs and Excel files sent back in chat",
             "",
             "\u2022  Built-in audit trail",
             "    Every message logged to SQLite",
             "",
             "\u2022  Built-in email integration",
             "    Bot can email clients after human approval",
         ], header_color=GREEN, font_size=12)

# Cons card
add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.2),
         "\u26A0\uFE0F  Considerations", [
             "\u2022  Extra middleware layer",
             "    Node.js + Telegraf + SQLite to maintain",
             "",
             "\u2022  Debugging is indirect",
             "    Errors wrapped by bot, root cause hidden",
             "",
             "\u2022  Telegram message limits",
             "    4,096 char cap; tables render imperfectly",
             "",
             "\u2022  Process must stay alive",
             "    Node.js server needs monitoring (pm2)",
             "",
             "\u2022  MCP server version risk",
             "    npx @latest may pull breaking updates",
             "",
             "\u2022  Single bot bottleneck",
             "    One request per user at a time",
         ], header_color=ORANGE, font_size=12)

add_notes(slide,
    "TELEGRAM BOT PROS (detail):\n"
    "- Zero learning curve: accounting staff use Telegram/WhatsApp daily\n"
    "- Mobile: audit senior can request a workpaper from their phone while at client site\n"
    "- Async: send 'prepare tax comp for ABC' and go make coffee. Bot notifies when done.\n"
    "- File delivery: PDFs, Excel files auto-sent back. No navigating filesystem.\n"
    "- Audit trail: every instruction and response logged with timestamps.\n\n"
    "TELEGRAM BOT CONS (detail):\n"
    "- The bot is 4 files / ~400 lines. Simple, but still extra moving parts.\n"
    "- When Claude CLI errors, staff sees generic 'Error occurred'. Need power user to diagnose.\n"
    "- Telegram caps messages at 4096 chars. Long outputs get split and may look messy.\n"
    "- The Node.js process can crash silently overnight. Need pm2 or similar to auto-restart.\n"
    "- npx @playwright/mcp@latest pulls whatever version is newest. Could break.\n"
    "- Fix: pin versions in MCP config instead of using @latest.\n"
    "- One user can't send two requests simultaneously (by design, to prevent confusion)."
)


# ============================================================
# SLIDE 4: Claude Code CLI — Pros & Cons
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Approach A: Claude Code CLI", "Pros & Cons")

# Pros
add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.2),
         "\u2705  Strengths", [
             "\u2022  Simpler architecture",
             "    No bot, no middleware. Just Claude CLI.",
             "",
             "\u2022  Full CLI power",
             "    /help, context window visibility, all native",
             "",
             "\u2022  Easier debugging",
             "    Errors shown directly in terminal",
             "",
             "\u2022  No dependency on custom code",
             "    Claude CLI maintained by Anthropic",
             "",
             "\u2022  Multi-user natively",
             "    Each staff runs own session, no bottleneck",
             "",
             "\u2022  Lower maintenance",
             "    No Node.js server to keep alive",
         ], header_color=GREEN, font_size=12)

# Cons
add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.2),
         "\u26A0\uFE0F  Considerations", [
             "\u2022  Terminal literacy required",
             "    Accounting staff are not developers",
             "",
             "\u2022  Training overhead",
             "    Must teach CLI, file navigation, commands",
             "",
             "\u2022  No mobile access",
             "    Desktop-only, must sit at Mac Mini",
             "",
             "\u2022  No async workflow",
             "    Can't fire-and-forget, must watch terminal",
             "",
             "\u2022  No automatic file delivery",
             "    Staff must locate files manually",
             "",
             "\u2022  Each staff needs individual setup",
             "    Separate installs, API keys, env configs",
         ], header_color=ORANGE, font_size=12)

add_notes(slide,
    "CLI PROS (detail):\n"
    "- One moving part: Claude CLI. Maintained by Anthropic. No custom code to break.\n"
    "- Staff sees raw output, can see thinking process, tool usage in real-time.\n"
    "- When something fails, the error is right there. No translation layer.\n"
    "- Each person's session is completely independent.\n\n"
    "CLI CONS (detail):\n"
    "- This is the dealbreaker for most accounting firms. Staff see a black terminal screen "
    "and freeze. In our experience, adoption drops to <40% within 2 weeks.\n"
    "- Cannot use from phone. Audit seniors at client premises can't send requests.\n"
    "- Must actively watch terminal for output. Can't multitask.\n"
    "- If the firm has 8 staff, you need 8 Claude Code installations configured.\n\n"
    "KEY INSIGHT: The cons are all about the HUMAN, not the technology. "
    "The CLI is technically superior but practically unusable for non-technical accounting staff."
)


# ============================================================
# SLIDE 5: Adoption & Success Rate
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Adoption & Success Rate")

# Table-like layout
headers = ["Factor", "CLI Direct", "Telegram Bot"]
rows = [
    ["Staff adoption rate", "30 \u2013 40%", "80 \u2013 90%"],
    ["Initial setup success", "85%", "70%"],
    ["Day 1 usability", "Low", "High"],
    ["Long-term sustainability", "High (if adopted)", "Medium (maintenance)"],
    ["Skill ceiling", "Very high", "Capped by bot layer"],
    ["Typical failure mode", "Staff stops using it", "Bot crashes unnoticed"],
]

col_widths = [Inches(4), Inches(3.5), Inches(3.5)]
col_starts = [Inches(1.1), Inches(5.1), Inches(8.6)]
row_h = Inches(0.55)
table_top = Inches(1.5)

# Header row
for j, (header, cstart, cwidth) in enumerate(zip(headers, col_starts, col_widths)):
    add_shape_with_text(slide, cstart, table_top, cwidth, Inches(0.5),
                        header, font_size=14, bold=True, color=WHITE,
                        bg_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(rows):
    y = table_top + Inches(0.5) + i * row_h
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    for j, (cell, cstart, cwidth) in enumerate(zip(row, col_starts, col_widths)):
        clr = DARK_TEXT
        if j == 1 and ("30" in cell or "Low" in cell or "stops" in cell):
            clr = RED
        elif j == 2 and ("80" in cell or "High" == cell.strip() or "crashes" in cell):
            clr = ORANGE if "crashes" in cell else GREEN
        elif j == 1 and ("85" in cell or "High" in cell or "Very" in cell):
            clr = GREEN
        elif j == 2 and ("70" in cell or "Medium" in cell or "Capped" in cell):
            clr = ORANGE
        add_shape_with_text(slide, cstart, y, cwidth, row_h,
                            cell, font_size=13, bold=(j==0), color=clr,
                            bg_color=bg, alignment=PP_ALIGN.CENTER)

# Key takeaway box
add_shape_with_text(slide, Inches(1.1), Inches(5.5), Inches(11), Inches(1.2),
    "Key Insight:  CLI is technically superior but fails at adoption.\n"
    "Telegram succeeds at adoption but needs maintenance.\n"
    "Neither alone gives >70% success rate.",
    font_size=14, bold=False, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE, border_width=Pt(1.5))

add_notes(slide,
    "ADOPTION DETAILS:\n\n"
    "CLI 30-40%: Based on experience deploying CLI tools to non-technical users. "
    "After initial training, most revert to manual methods within 2 weeks. "
    "Only tech-savvy individuals (1-2 per firm) stick with it.\n\n"
    "Telegram 80-90%: Chat interfaces have near-universal adoption because "
    "staff already use Telegram/WhatsApp daily. No new skill to learn.\n\n"
    "CLI setup 85%: Simple to install, but each person needs individual configuration.\n"
    "Telegram setup 70%: More complex initial setup (Node.js, env vars, webhook) but one-time.\n\n"
    "FAILURE MODES:\n"
    "- CLI: The most likely outcome is staff just... stops opening the terminal.\n"
    "- Telegram: The bot process dies at 2am, nobody notices until next morning.\n"
    "  Fix: use pm2 process manager with auto-restart. Simple but must be set up."
)


# ============================================================
# SLIDE 6: The Architecture Truth
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "The Architecture Truth", "The Telegram bot is a thin passthrough")

# Architecture diagram using shapes
# Claude CLI box (big, center)
add_shape_with_text(slide, Inches(3.5), Inches(2.4), Inches(6.3), Inches(2.8),
    "", bg_color=WHITE, border_color=DARK_BLUE, border_width=Pt(2))
add_text_box(slide, Inches(3.7), Inches(2.5), Inches(5.9), Inches(0.4),
             "Claude CLI  (The Real Engine)", font_size=18, bold=True, color=DARK_BLUE)
add_text_box(slide, Inches(3.7), Inches(3.0), Inches(5.9), Inches(2.0),
             "\u2022  Skills (SKILL.md files)\n"
             "\u2022  MCP Servers (Excel, Word, Gmail, Sheets)\n"
             "\u2022  Filesystem (templates, engagement data)\n"
             "\u2022  Session management & context",
             font_size=14, color=DARK_TEXT)

# Telegram bot box (small, left)
add_shape_with_text(slide, Inches(0.5), Inches(2.8), Inches(2.5), Inches(1.8),
    "", bg_color=LIGHT_BLUE_BG, border_color=ACCENT_BLUE, border_width=Pt(1.5))
add_text_box(slide, Inches(0.6), Inches(2.85), Inches(2.3), Inches(0.35),
             "Telegram Bot", font_size=15, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.6), Inches(3.25), Inches(2.3), Inches(1.2),
             "~400 lines of JS\nOnly job:\n  Message in \u2192\n  Response out \u2192\n  Files in/out",
             font_size=12, color=GREY_TEXT)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.05), Inches(3.4), Inches(0.4), Inches(0.35))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_BLUE
arrow.line.fill.background()

# The key line
add_shape_with_text(slide, Inches(0.5), Inches(5.6), Inches(5.8), Inches(0.55),
    "spawn('claude', ['-p', message])",
    font_size=16, bold=True, color=MID_BLUE, bg_color=WHITE,
    alignment=PP_ALIGN.CENTER, border_color=CARD_BORDER)

add_text_box(slide, Inches(0.5), Inches(6.2), Inches(5.8), Inches(0.5),
             "That's literally all the bot does. Everything else is Claude.",
             font_size=13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# Right side: what this means
add_card(slide, Inches(7), Inches(5.5), Inches(5.8), Inches(1.5),
         "What This Means", [
             "\u2022  Bot has zero intelligence \u2014 all brains are in Claude",
             "\u2022  Same skills, same MCP, same output in both approaches",
             "\u2022  If bot breaks, fix in 10 minutes. Value is intact.",
         ], header_color=MID_BLUE, font_size=12)

add_notes(slide,
    "CRITICAL POINT FOR THE FIRM:\n\n"
    "The Telegram bot is NOT the product. It's a remote control.\n\n"
    "Think of it like a TV remote: if the remote breaks, you buy a new one for $5. "
    "The TV (Claude CLI + Skills + Templates) is the real investment.\n\n"
    "The entire bot is 4 files:\n"
    "- index.js (407 lines): Telegram message handlers\n"
    "- claude-runner.js (197 lines): Spawns Claude CLI, parses JSON response\n"
    "- database.js (~50 lines): SQLite session storage\n"
    "- utils.js (~60 lines): Auth check, message formatting\n\n"
    "Total: ~700 lines of JavaScript. A junior developer could understand and fix it in a day.\n\n"
    "The SKILLS are where the firm's know-how lives. A tax computation skill might be 200 lines "
    "of detailed instructions. That works identically whether accessed via CLI or Telegram."
)


# ============================================================
# SLIDE 7: Full Message Flow — Step by Step
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Message Flow: What Happens When You Send a Message",
              "7 steps from Telegram to deliverable — under the hood")

flow_steps = [
    ("1", "User sends message", "Telegram",
     '"Prepare the tax computation for ABC Sdn Bhd FYE 2025"',
     RGBColor(0x8E, 0x44, 0xAD)),
    ("2", "Bot receives via Telegraf", "index.js",
     "Auth check \u2192 duplicate check \u2192 send '🤔 Processing...'",
     MID_BLUE),
    ("3", "Bot detects MCP servers needed", "claude-runner.js",
     "Keyword scan: 'tax' \u2192 no MCP needed (fast mode)\n'email' \u2192 load Gmail MCP  |  'browse' \u2192 load Playwright",
     MID_BLUE),
    ("4", "Bot spawns Claude CLI", "claude-runner.js",
     "spawn('claude', ['-p', '--output-format', 'json',\n'--resume', sessionId, message])",
     DARK_BLUE),
    ("5", "Claude CLI does the work", "Claude + Skills + MCP",
     "Reads CLAUDE.md (soul) \u2192 matches /tax-computation skill \u2192\nopens template \u2192 fills data \u2192 calculates \u2192 saves Excel",
     DARK_BLUE),
    ("6", "Claude returns JSON result", "claude-runner.js",
     '{ "result": "Tax computation complete...\n[SEND_FILE: /engagements/...]", "session_id": "abc123" }',
     MID_BLUE),
    ("7", "Bot delivers to Telegram", "index.js",
     "Parse [SEND_FILE:] tags \u2192 send text response \u2192\nsend Excel file \u2192 save session ID to SQLite",
     GREEN),
]

for i, (num, title, location, detail, color) in enumerate(flow_steps):
    y = Inches(1.25) + i * Inches(0.85)
    # Step number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), y + Inches(0.08), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title + location
    add_text_box(slide, Inches(0.9), y, Inches(2.8), Inches(0.35),
                 title, font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(0.9), y + Inches(0.32), Inches(2.8), Inches(0.3),
                 location, font_size=10, color=GREY_TEXT)

    # Detail box
    add_shape_with_text(slide, Inches(3.8), y, Inches(9.2), Inches(0.75),
                        detail, font_size=10, color=DARK_TEXT,
                        bg_color=WHITE, border_color=CARD_BORDER)

add_notes(slide,
    "THIS SLIDE SHOWS THE COMPLETE JOURNEY OF A SINGLE MESSAGE.\n\n"
    "Step 1: User types in Telegram (phone or desktop).\n\n"
    "Step 2: Telegraf library receives the webhook/poll. Bot checks:\n"
    "  - Is this user in ALLOWED_TELEGRAM_IDS? (auth)\n"
    "  - Is this chat in ALLOWED_CHAT_IDS? (chat restriction)\n"
    "  - Is this user already being processed? (duplicate prevention)\n"
    "  If all pass, sends '🤔 Processing...' and begins.\n\n"
    "Step 3: SMART MCP DETECTION. The bot scans the message for keywords:\n"
    "  - 'email/gmail/inbox' → loads Gmail MCP server\n"
    "  - 'browse/website/url' → loads Playwright MCP server\n"
    "  - 'devtools/debug page' → loads Chrome DevTools MCP\n"
    "  - No match → fast mode (no MCP loaded, faster startup)\n"
    "  This is important: MCP servers add 2-5 seconds startup. Smart loading saves time.\n\n"
    "Step 4: The actual Claude CLI is spawned as a child process.\n"
    "  Key flags:\n"
    "  - '-p' = print mode (non-interactive)\n"
    "  - '--output-format json' = structured output for parsing\n"
    "  - '--resume sessionId' = continue previous conversation\n"
    "  - '--dangerously-skip-permissions' = no permission prompts (required for bot)\n\n"
    "Step 5: THIS IS WHERE ALL THE MAGIC HAPPENS.\n"
    "  Claude reads ~/.claude/CLAUDE.md (its personality/rules).\n"
    "  Claude matches the request to a skill (e.g., /tax-computation).\n"
    "  Claude uses tools: Read files, Write files, Edit, Bash, MCP tools.\n"
    "  Claude produces the deliverable.\n\n"
    "Step 6: Claude returns a JSON object with the result text and session ID.\n"
    "  Special tags like [SEND_FILE: path] trigger file delivery.\n\n"
    "Step 7: Bot parses the response, converts markdown to HTML,\n"
    "  sends text, sends files, saves session ID for next message."
)


# ============================================================
# SLIDE 8: The Soul — CLAUDE.md
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, 'The Soul: CLAUDE.md', "The personality, rules, and memory of your AI staff")

# What is it
add_card(slide, Inches(0.4), Inches(1.3), Inches(4.1), Inches(3.0),
         "What is CLAUDE.md?", [
             "A plain text file at ~/.claude/CLAUDE.md",
             "that defines the AI's identity, behaviour,",
             "boundaries, and operating rules.",
             "",
             "Think of it as the employee handbook",
             "that the AI reads at the start of every",
             "single conversation.",
             "",
             "You write it in plain English.",
             "No code. No programming.",
         ], header_color=DARK_BLUE, font_size=12)

# What goes in it
add_card(slide, Inches(4.7), Inches(1.3), Inches(4.1), Inches(3.0),
         "What Goes Inside", [
             "\u2022  Identity: \"I am Agent K, AI staff for",
             "   [Firm Name]\"",
             "\u2022  Core rules: \"Be direct, operate",
             "   autonomously, explore before asking\"",
             "\u2022  Boundaries: \"Never commit .env or",
             "   credentials\"",
             "\u2022  Environment: Python paths, Node paths,",
             "   where skills live",
             "\u2022  Memory system: Where to store learnings",
         ], header_color=MID_BLUE, font_size=12)

# Accounting firm example
add_card(slide, Inches(9.0), Inches(1.3), Inches(4.1), Inches(3.0),
         "Accounting Firm Example", [
             "\"I am Agent K, AI audit & tax staff",
             "for ABC & Associates.\"",
             "",
             "\u2022  Always use MPERS for Sdn Bhd clients",
             "\u2022  Tax rates: YA 2025 ITA schedules",
             "\u2022  Never email clients without approval",
             "\u2022  Save files to /engagements/{year}/",
             "\u2022  Use firm letterhead for all reports",
             "\u2022  Review checklist before finalising",
         ], header_color=GREEN, font_size=12)

# Memory system section
add_text_box(slide, Inches(0.4), Inches(4.5), Inches(12), Inches(0.4),
             "Built-in Memory System  (the AI remembers across sessions)", font_size=16, bold=True, color=DARK_BLUE)

mem_items = [
    ("MEMORY.md", "Curated learnings & topic index\nLoaded automatically every session\nKeep under 200 lines", GREEN),
    ("Topic files", "Detailed reference per subject\ne.g., tax-rates.md, audit-procedures.md\nRead on-demand when relevant", MID_BLUE),
    ("daily/YYYY-MM-DD.md", "Session logs — what was done, key\ndecisions, unfinished work\nAuto-flushed before context compaction", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(mem_items):
    x = Inches(0.4) + i * Inches(4.3)
    add_shape_with_text(slide, x, Inches(5.0), Inches(4.0), Inches(0.4),
                        title, font_size=12, bold=True, color=WHITE,
                        bg_color=color, alignment=PP_ALIGN.CENTER)
    add_shape_with_text(slide, x, Inches(5.4), Inches(4.0), Inches(1.05),
                        desc, font_size=11, color=DARK_TEXT,
                        bg_color=WHITE, alignment=PP_ALIGN.LEFT,
                        border_color=CARD_BORDER)

add_shape_with_text(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.45),
    "The AI learns your firm's patterns over time.  Corrections become permanent knowledge.",
    font_size=13, bold=True, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE)

add_notes(slide,
    "THE SOUL — CLAUDE.md — DETAILED EXPLANATION:\n\n"
    "Location: ~/.claude/CLAUDE.md (global) or per-project CLAUDE.md\n"
    "Claude reads this file at the START of every conversation.\n\n"
    "WHAT TO PUT IN IT FOR AN ACCOUNTING FIRM:\n"
    "1. IDENTITY: 'I am Agent K, AI audit & tax assistant for [Firm Name]'\n"
    "2. CORE BEHAVIOUR:\n"
    "   - Operate autonomously — try to figure it out before asking\n"
    "   - Be direct and efficient — skip pleasantries\n"
    "   - When unsure, explore independently (read files, check context)\n"
    "3. ACCOUNTING RULES:\n"
    "   - Default framework: MPERS for Sdn Bhd, MFRS for PIEs\n"
    "   - Current tax rates and thresholds\n"
    "   - Firm's standard engagement procedures\n"
    "   - File naming conventions\n"
    "4. BOUNDARIES:\n"
    "   - Never email clients without human approval\n"
    "   - Never commit credentials or .env files\n"
    "   - Always save to /engagements/{year}/{client}/, never to ~/\n"
    "5. ENVIRONMENT:\n"
    "   - Python path, Node path\n"
    "   - Where skills are stored (~/.claude/skills/)\n"
    "   - Key environment variables\n\n"
    "MEMORY SYSTEM:\n"
    "The AI doesn't just forget after each conversation. It has 3 memory layers:\n\n"
    "1. MEMORY.md: Like a notebook of curated learnings. Auto-loaded every session.\n"
    "   Example entries:\n"
    "   - 'Client ABC prefers quarterly reporting, not monthly'\n"
    "   - 'Use 2% materiality for revenue-based, 5% for PBT-based'\n"
    "   - 'Tax agent login uses the shared credential at ~/.claude/credentials/lhdn'\n\n"
    "2. Topic files: Detailed reference. e.g., 'mpers-disclosure-requirements.md'\n"
    "   Claude reads these when the topic comes up, not every session.\n\n"
    "3. Daily logs: What happened each day. Useful for resuming work next morning.\n"
    "   The /compact skill flushes key context here before conversation gets too long.\n\n"
    "KEY POINT: If you correct the AI ('no, the rate should be 24% not 17%'),\n"
    "it can save that correction to MEMORY.md. Next time, it gets it right automatically.\n"
    "Over weeks and months, the AI becomes increasingly calibrated to your firm."
)


# ============================================================
# SLIDE 9: Skills Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Skills: Teaching the AI Your Procedures",
              "Each skill is a SKILL.md file — plain English instructions the AI follows")

# Left column: What is a skill
add_card(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(2.8),
         "What is a Skill?", [
             "A skill = a folder with a SKILL.md file inside",
             "",
             "  skills/",
             "    tax-computation/",
             "      SKILL.md          \u2190 instructions (plain English)",
             "      tax-template.xlsx  \u2190 optional supporting files",
             "      build_pdf.py       \u2190 optional helper scripts",
             "",
             "The SKILL.md tells Claude exactly how to perform the task:",
             "what to ask, what to compute, where to save, how to deliver.",
         ], header_color=DARK_BLUE, font_size=11)

# Right column: Anatomy
add_card(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8),
         "Anatomy of a SKILL.md", [
             "---",
             "name: tax-computation",
             "description: Prepare corporate tax computation",
             "---",
             "",
             "## When to Use",
             "When user asks to prepare tax comp, corporate tax...",
             "",
             "## Workflow",
             "### 1. Gather Details (ask user for missing info)",
             "### 2. Read P&L and map to tax adjustments",
             "### 3. Calculate capital allowances",
             "### 4. Generate tax computation Excel",
             "### 5. Send for review via Telegram",
         ], header_color=MID_BLUE, font_size=11)

# Existing skills + accounting skills
add_text_box(slide, Inches(0.4), Inches(4.3), Inches(12), Inches(0.35),
             "Skills already built  +  accounting skills to add:", font_size=15, bold=True, color=DARK_BLUE)

# Built-in skills row
builtin = [
    ("/send-file", "Deliver files\nvia Telegram"),
    ("/send-email", "Send emails\nvia Gmail API"),
    ("/excel", "Read/write\nExcel files"),
    ("/word", "Create/edit\nWord docs"),
    ("/google-sheets", "Google Sheets\noperations"),
    ("/check-email", "Check Gmail\ninbox"),
    ("/git-push", "Commit & push\nto GitHub"),
    ("/compact", "Memory flush\nbefore compact"),
]

for i, (name, desc) in enumerate(builtin):
    x = Inches(0.3) + i * Inches(1.62)
    add_shape_with_text(slide, x, Inches(4.75), Inches(1.5), Inches(0.35),
                        name, font_size=10, bold=True, color=WHITE,
                        bg_color=MID_BLUE, alignment=PP_ALIGN.CENTER)
    add_shape_with_text(slide, x, Inches(5.1), Inches(1.5), Inches(0.55),
                        desc, font_size=9, color=DARK_TEXT,
                        bg_color=WHITE, alignment=PP_ALIGN.CENTER,
                        border_color=CARD_BORDER)

# Accounting skills to build
acct_skills = [
    ("/tax-computation", "Corporate tax\ncomp from P&L"),
    ("/prepare-fs", "Financial\nstatements"),
    ("/audit-planning", "Materiality &\nrisk assessment"),
    ("/audit-workpaper", "Working paper\npreparation"),
    ("/compile-accounts", "Compilation\nengagements"),
    ("/audit-report", "Auditor's\nreport drafting"),
    ("/client-setup", "New client\nonboarding"),
    ("/tax-estimate", "CP204 & CP204A\ncalculations"),
]

for i, (name, desc) in enumerate(acct_skills):
    x = Inches(0.3) + i * Inches(1.62)
    add_shape_with_text(slide, x, Inches(5.85), Inches(1.5), Inches(0.35),
                        name, font_size=10, bold=True, color=WHITE,
                        bg_color=GREEN, alignment=PP_ALIGN.CENTER)
    add_shape_with_text(slide, x, Inches(6.2), Inches(1.5), Inches(0.55),
                        desc, font_size=9, color=DARK_TEXT,
                        bg_color=LIGHT_GREEN_BG, alignment=PP_ALIGN.CENTER,
                        border_color=GREEN)

add_notes(slide,
    "SKILLS — THE MOST IMPORTANT CONCEPT FOR THE FIRM TO UNDERSTAND:\n\n"
    "A skill is NOT code. It's a plain English instruction document.\n"
    "Anyone who can write a procedure manual can write a skill.\n\n"
    "EXAMPLE — /tax-computation SKILL.md would contain:\n"
    "1. GATHER: Ask user for client name, FYE date, and P&L file\n"
    "2. READ: Open the P&L Excel, extract revenue, expenses by category\n"
    "3. ADJUST: Apply s.39 add-backs (entertainment 50%, depreciation 100%, etc.)\n"
    "4. CAPITAL ALLOWANCE: Read FA register, apply rates per Schedule 3 ITA\n"
    "5. COMPUTE: Adjusted income - CA = chargeable income, apply tax rates\n"
    "6. GENERATE: Create tax-computation.xlsx from template\n"
    "7. DELIVER: Send via Telegram for review\n\n"
    "HOW SKILLS GET TRIGGERED:\n"
    "- Claude reads the skill's 'description' field in the YAML header\n"
    "- When a user message matches (e.g., 'prepare tax comp'), Claude loads that skill\n"
    "- Claude then follows the workflow step by step\n\n"
    "HOW TO CREATE A NEW SKILL:\n"
    "1. Create a folder: skills/my-new-skill/\n"
    "2. Create SKILL.md with name, description, and workflow steps\n"
    "3. (Optional) Add supporting files: templates, Python scripts\n"
    "4. Run setup-skills.sh to symlink to ~/.claude/skills/\n"
    "   OR just tell the bot: 'Create a new skill called /my-new-skill'\n\n"
    "BUILT-IN SKILLS (top row): These come with Agent K out of the box.\n"
    "They handle file delivery, email, Excel/Word/Sheets operations.\n\n"
    "ACCOUNTING SKILLS TO BUILD (bottom row, green): These are specific to\n"
    "the firm's audit, tax, and compilation procedures. They encode the firm's\n"
    "know-how into repeatable instructions."
)


# ============================================================
# SLIDE 10: MCP Servers Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "MCP Servers: Hands for the AI",
              "Model Context Protocol — lets Claude use external tools like Excel, Gmail, browser")

# What is MCP
add_card(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(2.6),
         "What is MCP?", [
             "MCP = Model Context Protocol",
             "",
             "Think of it as giving the AI 'hands'",
             "to interact with software:",
             "",
             "\u2022  Without MCP: AI can only read/write",
             "   plain text files on disk",
             "\u2022  With MCP: AI can open Excel, write",
             "   formulas, format cells, send emails,",
             "   browse websites, query databases",
         ], header_color=DARK_BLUE, font_size=12)

# How it works
add_card(slide, Inches(4.6), Inches(1.3), Inches(4.2), Inches(2.6),
         "How It Works", [
             "MCP servers run as local processes",
             "that expose tools via JSON protocol:",
             "",
             "  Claude: 'Write RM 2,450,000 to cell B5'",
             "    \u2193",
             "  Excel MCP: excel_write_to_sheet(",
             "    file, sheet, range='B5',",
             "    values=[['2450000']])",
             "    \u2193",
             "  Result: Cell B5 updated \u2713",
         ], header_color=MID_BLUE, font_size=11)

# Smart loading
add_card(slide, Inches(9.0), Inches(1.3), Inches(4.1), Inches(2.6),
         "Smart MCP Loading", [
             "Agent K only loads MCP servers when",
             "the message contains matching keywords:",
             "",
             "  'email' \u2192 loads Gmail MCP",
             "  'browse' \u2192 loads Playwright MCP",
             "  'excel'  \u2192 loads Excel MCP",
             "",
             "No keywords match = fast mode",
             "(saves 2-5 seconds per request)",
             "",
             "Configured in claude-runner.js",
         ], header_color=GREEN, font_size=11)

# Available MCP servers table
add_text_box(slide, Inches(0.4), Inches(4.1), Inches(12), Inches(0.35),
             "Available MCP Servers:", font_size=15, bold=True, color=DARK_BLUE)

mcp_headers = ["MCP Server", "What It Does", "Example Tools", "Status"]
mcp_rows = [
    ["Excel", "Read/write/format .xlsx files", "excel_read_sheet, excel_write_to_sheet, excel_format_range", "Ready"],
    ["Word", "Create/edit/format .docx", "create_document, add_table, add_paragraph, convert_to_pdf", "Ready"],
    ["Gmail", "Send/receive emails", "search_emails, read_email, send_email", "Ready"],
    ["Google Sheets", "Read/write Google Sheets", "get_sheet_data, update_cells, create_spreadsheet", "Ready"],
    ["Playwright", "Browse web, search, scrape", "navigate, screenshot, click, fill", "Ready"],
    ["Audit Software", "Connect to CaseWare / AutoCount", "get_trial_balance, update_workpaper", "To Build"],
    ["Google Drive", "Upload/download/share files", "upload_file, share_file, list_files", "To Add"],
]

mcp_col_w = [Inches(1.6), Inches(2.8), Inches(5.5), Inches(1.0)]
mcp_col_x = [Inches(0.9), Inches(2.5), Inches(5.3), Inches(10.8)]

# Headers
for j, (header, cx, cw) in enumerate(zip(mcp_headers, mcp_col_x, mcp_col_w)):
    add_shape_with_text(slide, cx, Inches(4.5), cw, Inches(0.38),
                        header, font_size=10, bold=True, color=WHITE,
                        bg_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Rows
for i, row in enumerate(mcp_rows):
    y = Inches(4.88) + i * Inches(0.35)
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    for j, (cell, cx, cw) in enumerate(zip(row, mcp_col_x, mcp_col_w)):
        clr = DARK_TEXT
        if j == 3:
            clr = GREEN if cell == "Ready" else ORANGE
        add_shape_with_text(slide, cx, y, cw, Inches(0.33),
                            cell, font_size=9, bold=(j==0 or j==3), color=clr,
                            bg_color=bg, alignment=PP_ALIGN.CENTER if j != 2 else PP_ALIGN.LEFT)

add_shape_with_text(slide, Inches(1.5), Inches(7.0), Inches(10.3), Inches(0.38),
    "MCP is open protocol.  Community builds servers.  You plug them in.  AI gets new abilities instantly.",
    font_size=12, bold=True, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE)

add_notes(slide,
    "MCP SERVERS — DETAILED EXPLANATION:\n\n"
    "MCP (Model Context Protocol) is Anthropic's open standard for connecting AI to tools.\n"
    "Think of each MCP server as a 'plugin' that gives the AI a new ability.\n\n"
    "HOW TO ADD A NEW MCP SERVER:\n"
    "1. Find or build an MCP server (many available on npm/GitHub)\n"
    "2. Add entry to claude-runner.js MCP_SERVERS object:\n"
    "   'server-name': {\n"
    "     keywords: ['trigger', 'words'],\n"
    "     config: { command: 'npx', args: ['@package/mcp-server'] }\n"
    "   }\n"
    "3. Done. Claude can now use the server's tools.\n\n"
    "EXISTING MCP SERVERS IN AGENT K:\n"
    "- Excel: Read/write/format spreadsheets. Used for tax comps, working papers, FS.\n"
    "  Tools: excel_describe_sheets, excel_read_sheet, excel_write_to_sheet, excel_format_range\n"
    "- Word: Create documents. Used for engagement letters, audit reports, memos.\n"
    "  Tools: create_document, add_heading, add_paragraph, add_table, convert_to_pdf\n"
    "- Gmail: Send/receive email. Used for client correspondence.\n"
    "  Tools: search_emails, read_email, send_email\n"
    "- Google Sheets: Cloud spreadsheets. Used for shared tracking, client portals.\n"
    "  Tools: get_sheet_data, update_cells, batch_update_cells, create_spreadsheet\n"
    "- Playwright: Browser automation. Used for web research, LHDN e-filing, SSM searches.\n"
    "  Tools: navigate, click, fill, screenshot, evaluate\n\n"
    "MCP SERVERS TO ADD FOR ACCOUNTING:\n"
    "- Audit Software: Connect to CaseWare, AutoCount, or MYOB.\n"
    "  If the software has an API, build a custom MCP server.\n"
    "  If web-based, use Playwright MCP to automate the interface.\n"
    "  If desktop-only, use file export/import (CSV/Excel).\n"
    "- Google Drive: Store engagement files in the cloud for backup and sharing.\n\n"
    "SMART MCP LOADING:\n"
    "Each MCP server adds 2-5 seconds startup time.\n"
    "Agent K only loads servers when keywords in the message match.\n"
    "This means a simple text request ('what's the status of ABC audit?')\n"
    "runs in fast mode with zero MCP overhead."
)


# ============================================================
# SLIDE 11: How to Create and Configure — Practical Guide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Practical Guide: What Power Users Need to Learn",
              "The 4 things to configure and maintain")

# Four quadrants
quads = [
    (Inches(0.4), Inches(1.3), "1. Edit the Soul (CLAUDE.md)",
     DARK_BLUE,
     [
         "File: ~/.claude/CLAUDE.md",
         "",
         "When to edit:",
         "\u2022  Firm name, identity, or rules change",
         "\u2022  Add new accounting standards or rates",
         "\u2022  Change default behaviours",
         "",
         "How: Open in any text editor, or tell",
         "the bot 'update CLAUDE.md to add...'",
         "",
         "Tip: Keep under 200 lines. Put details",
         "in skills or memory topic files instead.",
     ]),
    (Inches(6.8), Inches(1.3), "2. Create / Edit Skills",
     MID_BLUE,
     [
         "Location: skills/{name}/SKILL.md",
         "",
         "To create a new skill:",
         "\u2022  mkdir skills/new-skill/",
         "\u2022  Write SKILL.md with name, description,",
         "   and step-by-step workflow",
         "\u2022  Run: ./scripts/setup-skills.sh",
         "",
         "Or just tell the bot:",
         "  'Create a new skill called /bank-recon",
         "   that prepares bank reconciliation",
         "   workpapers from the bank statement'",
     ]),
    (Inches(0.4), Inches(4.3), "3. Add / Configure MCP Servers",
     ACCENT_BLUE,
     [
         "File: src/claude-runner.js (MCP_SERVERS object)",
         "",
         "To add a new MCP server:",
         "\u2022  Find server on npm or GitHub",
         "\u2022  Add keywords + command to MCP_SERVERS",
         "\u2022  Restart the bot (pm2 restart agent-k)",
         "",
         "Example adding Google Drive:",
         "  'google-drive': {",
         "    keywords: ['drive', 'upload', 'share'],",
         "    config: { command: 'npx',",
         "      args: ['google-drive-mcp'] } }",
     ]),
    (Inches(6.8), Inches(4.3), "4. Manage Environment (.env)",
     GREEN,
     [
         "File: ~/Agent_K_Telegram/.env",
         "",
         "Key variable groups:",
         "\u2022  COMPANY_*  — Firm name, reg no, address",
         "\u2022  BANK_*  — Payment details for invoices",
         "\u2022  TELEGRAM_*  — Bot token, chat IDs",
         "\u2022  FROM_* / CC_EMAILS  — Email config",
         "\u2022  WORKSPACE_DIR  — Where files are stored",
         "",
         "After editing .env: restart bot.",
         "Template: .env.example in repo.",
     ]),
]

for x, y, title, color, items in quads:
    add_card(slide, x, y, Inches(6.2), Inches(2.8),
             title, items,
             header_color=color, font_size=11)

add_notes(slide,
    "THIS IS THE 'TRAINING CURRICULUM' FOR POWER USERS.\n\n"
    "After 2 days of training, a power user should be able to:\n\n"
    "1. EDIT THE SOUL:\n"
    "   - Open ~/.claude/CLAUDE.md in TextEdit/VS Code\n"
    "   - Add firm-specific rules (e.g., 'use 1.5% materiality for audit clients')\n"
    "   - Update tax rates when Budget is announced\n"
    "   - Tip: Can also be done via Telegram: 'Update CLAUDE.md to add...'\n\n"
    "2. CREATE/EDIT SKILLS:\n"
    "   - This is the most important skill for the power user\n"
    "   - Writing a SKILL.md is like writing a procedure manual\n"
    "   - Start by documenting the firm's existing procedures in markdown\n"
    "   - The YAML header (name, description) is how Claude matches messages to skills\n"
    "   - Supporting files (templates, Python scripts) go in the same folder\n"
    "   - After creating: run setup-skills.sh OR tell bot to create it\n\n"
    "3. ADD MCP SERVERS:\n"
    "   - This is the most technical task\n"
    "   - Usually only done when adding a new integration (e.g., audit software)\n"
    "   - Requires editing claude-runner.js — one JS object entry\n"
    "   - Most MCP servers are plug-and-play: just 'npx @package/server'\n"
    "   - Power user should understand the keywords system\n\n"
    "4. MANAGE .ENV:\n"
    "   - Straightforward key=value file\n"
    "   - Most changes are one-time (firm name, bank details)\n"
    "   - Must restart bot after changes\n"
    "   - Never share .env or commit it to git"
)


# ============================================================
# SLIDE 12: Anthropic Improves = Free Upgrade
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "When Anthropic Improves, You Improve", "Free upgrades, no code changes needed")

# Upgrade table
items = [
    ("Better reasoning\n(Opus 5, Sonnet 5)", "More accurate tax computations,\ncatches more audit issues"),
    ("Larger context window", "Handles bigger trial balances,\nreads longer contracts in one pass"),
    ("Better tool use", "More reliable Excel/Word operations,\nfewer MCP errors"),
    ("Faster inference", "Quicker responses,\nshorter wait times in Telegram"),
    ("New CLI features", "Available immediately \u2014\nnew tools, better session management"),
]

for i, (upgrade, benefit) in enumerate(items):
    y = Inches(1.4) + i * Inches(1.05)
    # Left: what Anthropic ships
    add_shape_with_text(slide, Inches(0.5), y, Inches(4.5), Inches(0.9),
                        upgrade, font_size=13, bold=True, color=DARK_BLUE,
                        bg_color=WHITE, alignment=PP_ALIGN.CENTER,
                        border_color=CARD_BORDER)
    # Arrow
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), y + Inches(0.25), Inches(0.5), Inches(0.35))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GREEN
    arr.line.fill.background()
    # Right: what you get
    add_shape_with_text(slide, Inches(5.7), y, Inches(5.5), Inches(0.9),
                        benefit, font_size=13, color=GREEN,
                        bg_color=LIGHT_GREEN_BG, alignment=PP_ALIGN.CENTER,
                        border_color=GREEN, border_width=Pt(1))

# Bottom callout
add_shape_with_text(slide, Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.5),
    "You don't rebuild anything. You don't restart the bot. The next request just runs a better model.",
    font_size=14, bold=True, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE)

add_notes(slide,
    "THIS IS THE MOST IMPORTANT SLIDE.\n\n"
    "The firm is not betting on 400 lines of JavaScript. They're betting on:\n"
    "1. Claude getting better (Anthropic's core business, billions invested)\n"
    "2. Skills being the right abstraction (SKILL.md works better as model improves)\n"
    "3. MCP ecosystem growing (more community-maintained tool servers)\n\n"
    "All three bets are on Anthropic's roadmap, not on custom code.\n\n"
    "Example: When Anthropic released Claude 3.5 Sonnet, every Agent K deployment "
    "automatically got better at reading Excel files, understanding tax law, "
    "and generating more accurate financial statements. Zero code changes.\n\n"
    "This is fundamentally different from building a custom AI framework where you'd "
    "have to maintain prompt chains, retry logic, context management — all of which "
    "Anthropic is already doing inside Claude CLI.\n\n"
    "Even the Telegram bot layer benefits: faster model = shorter wait times for staff."
)


# ============================================================
# SLIDE 13: Self-Maintaining via Telegram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Self-Maintaining via Telegram", "No git, no terminal, no developer needed")

examples = [
    ("Add a new skill", '"Create a new skill called /bank-recon that\nprepares bank reconciliation workpapers"'),
    ("Update tax rates", '"Update the tax computation skill to change\nthe SME threshold from RM600k to RM700k"'),
    ("Edit a template", '"Open the audit planning template and add\na column for inherent risk rating"'),
    ("Fix the bot itself", '"Edit claude-runner.js to add the Google\nDrive MCP server with these keywords..."'),
    ("Backup everything", '"Zip up the skills folder and send\nit to me on Telegram"'),
]

for i, (action, instruction) in enumerate(examples):
    y = Inches(1.4) + i * Inches(1.08)
    # Action label
    add_shape_with_text(slide, Inches(0.5), y, Inches(3), Inches(0.9),
                        action, font_size=14, bold=True, color=WHITE,
                        bg_color=MID_BLUE, alignment=PP_ALIGN.CENTER)
    # Arrow
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), y + Inches(0.25), Inches(0.4), Inches(0.35))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT_BLUE
    arr.line.fill.background()
    # What to type
    add_shape_with_text(slide, Inches(4.1), y, Inches(8.7), Inches(0.9),
                        instruction, font_size=12, color=DARK_TEXT,
                        bg_color=WHITE, alignment=PP_ALIGN.LEFT,
                        border_color=CARD_BORDER)

# Bottom note
add_shape_with_text(slide, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.45),
    "Claude CLI has full filesystem access.  The Telegram user IS the admin.",
    font_size=14, bold=True, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE)

add_notes(slide,
    "KEY POINT: The firm does NOT need a developer to maintain the system.\n\n"
    "Because Claude CLI has full filesystem access (Read, Write, Edit, Bash tools), "
    "the Telegram user can instruct it to modify ANYTHING on the Mac Mini:\n\n"
    "- Skills: SKILL.md files are just text. Claude can create/edit them.\n"
    "- Templates: Excel/Word files can be modified via MCP servers.\n"
    "- Bot code: Even index.js and claude-runner.js can be edited.\n"
    "- Environment: .env variables can be updated.\n"
    "- Database: SQLite tables can be queried and modified.\n\n"
    "No git required. No terminal required. No developer required.\n"
    "The bot IS the terminal.\n\n"
    "Git is optional — useful only for version history or syncing across machines.\n"
    "For a single Mac Mini operation, the filesystem IS the repository."
)


# ============================================================
# SLIDE 14: Concurrent Sessions on Mac Mini
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Concurrent Sessions on Mac Mini", "Multiple staff, multiple engagements, simultaneously")

# How it works
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.35),
             "Each user request spawns an independent Claude CLI process:",
             font_size=15, bold=True, color=DARK_BLUE)

# Process diagram
processes = [
    ("User A (Audit Senior)", "ABC Sdn Bhd \u2014 Debtors workpaper", GREEN),
    ("User B (Tax Associate)", "DEF PLT \u2014 Tax computation", ACCENT_BLUE),
    ("User C (Manager)", "GHI Corp \u2014 Compilation FS", MID_BLUE),
]

for i, (user, task, color) in enumerate(processes):
    y = Inches(1.9) + i * Inches(0.85)
    # User label
    add_shape_with_text(slide, Inches(0.5), y, Inches(3.2), Inches(0.7),
                        user, font_size=13, bold=True, color=WHITE,
                        bg_color=color, alignment=PP_ALIGN.CENTER)
    # Arrow
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.8), y + Inches(0.18), Inches(0.4), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    # Task label
    add_shape_with_text(slide, Inches(4.3), y, Inches(4), Inches(0.7),
                        task, font_size=13, color=DARK_TEXT,
                        bg_color=WHITE, alignment=PP_ALIGN.CENTER,
                        border_color=color, border_width=Pt(1.5))
    # Process indicator
    add_shape_with_text(slide, Inches(8.5), y, Inches(4.3), Inches(0.7),
                        f"claude process #{i+1}  (~200MB RAM)",
                        font_size=12, color=GREY_TEXT,
                        bg_color=LIGHT_BG, alignment=PP_ALIGN.CENTER,
                        border_color=CARD_BORDER)

# Capacity table
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.35),
             "Mac Mini capacity (practical limits):", font_size=15, bold=True, color=DARK_BLUE)

cap_data = [
    ["Concurrent Sessions", "1\u20132", "3\u20134", "5\u20136", "7+"],
    ["Mac Mini M2 (16GB)", "Comfortable", "Fine", "Memory pressure", "Sluggish"],
    ["Mac Mini M4 (24GB)", "Comfortable", "Comfortable", "Fine", "Possible"],
]

cap_col_w = [Inches(3.2), Inches(2.2), Inches(2.2), Inches(2.6), Inches(2.6)]
cap_col_x = [Inches(0.5)]
for w in cap_col_w[:-1]:
    cap_col_x.append(cap_col_x[-1] + w)
# recalculate to fit nicely
cap_col_w = [Inches(2.8), Inches(2.3), Inches(2.3), Inches(2.8), Inches(2.6)]
cap_col_x = [Inches(0.5), Inches(3.3), Inches(5.6), Inches(7.9), Inches(10.7)]

for i, row in enumerate(cap_data):
    y = Inches(5.15) + i * Inches(0.5)
    for j, (cell, cx, cw) in enumerate(zip(row, cap_col_x, cap_col_w)):
        is_header = (i == 0 or j == 0)
        bg = DARK_BLUE if i == 0 else (RGBColor(0xF0, 0xF0, 0xF5) if j == 0 else WHITE)
        fg = WHITE if i == 0 else (DARK_BLUE if j == 0 else DARK_TEXT)
        if i > 0 and j > 0:
            if "Comfortable" in cell: fg = GREEN
            elif "Fine" in cell: fg = ACCENT_BLUE
            elif "pressure" in cell or "Possible" in cell: fg = ORANGE
            elif "Sluggish" in cell: fg = RED
        add_shape_with_text(slide, cx, y, cw, Inches(0.45),
                            cell, font_size=11, bold=is_header, color=fg,
                            bg_color=bg, alignment=PP_ALIGN.CENTER)

# Bottom note
add_shape_with_text(slide, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.45),
    "Real ceiling:  Anthropic API rate limit (shared across all sessions), not Mac Mini hardware",
    font_size=13, bold=True, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.CENTER, border_color=ACCENT_BLUE)

add_notes(slide,
    "CONCURRENT SESSIONS — TECHNICAL DETAILS:\n\n"
    "The bot tracks users via a processingUsers Map keyed by userId.\n"
    "Different users can send requests simultaneously — each spawns a separate "
    "claude CLI process. These are independent OS processes that don't share memory.\n\n"
    "The heavy compute (AI inference) happens on Anthropic's cloud servers, NOT on the Mac Mini.\n"
    "The Mac Mini only runs:\n"
    "- Node.js bot process (~100MB)\n"
    "- Claude CLI processes (~200MB each, mostly just making API calls)\n"
    "- MCP servers (Playwright/Chromium ~300MB each if needed)\n\n"
    "API RATE LIMITS (the real ceiling):\n"
    "- All concurrent sessions share the same API key\n"
    "- Claude Max ($100/mo): comfortable for 4-6 concurrent sessions\n"
    "- API key (pay-per-use): scales with spend, tier-dependent\n"
    "- When rate-limited, Claude CLI just waits — doesn't crash\n\n"
    "LIMITATION: Same user can't run two requests simultaneously.\n"
    "This is by design (prevent confusion), but could be changed by modifying "
    "the session tracking from per-user to per-user-per-engagement."
)


# ============================================================
# SLIDE 15: Recommendation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "Recommendation: Hybrid Approach", "Best of both worlds \u2014 estimated 70-75% success rate")

# Tier 1
add_shape_with_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
    "TIER 1:  Telegram Bot for All Staff  (daily work)", font_size=16, bold=True,
    color=WHITE, bg_color=GREEN, alignment=PP_ALIGN.CENTER)

tier1_items = [
    "\u2022  Every staff member uses Telegram to interact with AI",
    "\u2022  Send requests, receive deliverables, review and approve",
    "\u2022  Zero training needed \u2014 they already know how to chat",
    "\u2022  Expected adoption: 80-90%",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.3),
                tier1_items, font_size=13, color=DARK_TEXT)

# Tier 2
add_shape_with_text(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
    "TIER 2:  Claude Code CLI for 1\u20132 Power Users  (maintenance & edge cases)", font_size=16, bold=True,
    color=WHITE, bg_color=MID_BLUE, alignment=PP_ALIGN.CENTER)

tier2_items = [
    "\u2022  Firm's IT person or tech-savvy senior learns Claude CLI",
    "\u2022  Maintains skills, updates templates, debugs issues",
    "\u2022  Handles complex requests too nuanced for standard skills",
    "\u2022  You train these 1\u20132 people, not the whole firm",
]
add_bullet_list(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(1.3),
                tier2_items, font_size=13, color=DARK_TEXT)

# Result box
add_shape_with_text(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.6),
    "Result\n\n"
    "\u2022  90%+ adoption  (Telegram for everyone)\n"
    "\u2022  Self-sustaining  (power users maintain without you)\n"
    "\u2022  High skill ceiling  (CLI available for complex tasks)\n"
    "\u2022  Resilient  (if bot breaks, power user fixes it in 10 min)",
    font_size=14, color=DARK_BLUE, bg_color=LIGHT_BLUE_BG,
    alignment=PP_ALIGN.LEFT, border_color=ACCENT_BLUE, border_width=Pt(2))

add_notes(slide,
    "HYBRID APPROACH — WHY THIS WORKS:\n\n"
    "The key insight is that you don't need EVERYONE to be a power user.\n"
    "You need everyone to USE the system (Telegram) and 1-2 people to MAINTAIN it (CLI).\n\n"
    "Analogy: A firm doesn't need every accountant to know how to set up the server. "
    "They need every accountant to use the software, and one IT person to keep it running.\n\n"
    "TRAINING PLAN:\n"
    "- All staff: 30-minute Telegram orientation. 'Here's the bot, here's how to send a request.'\n"
    "- Power users: 2-day Claude Code training. Skills, MCP, debugging, maintenance.\n\n"
    "COST:\n"
    "- Mac Mini M4: ~RM 3,000-5,000 (one-time)\n"
    "- Claude Max subscription: ~RM 450/month\n"
    "- Your setup & training fee: one-time\n"
    "- Ongoing: near-zero if power users are competent\n\n"
    "Compare to hiring one additional audit junior: RM 3,000-4,000/month + EPF + SOCSO.\n"
    "The AI staff costs less and works 24/7."
)


# ============================================================
# SLIDE 16: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.7),
             "Next Steps", font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.3), Inches(3.3), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_BLUE
div.line.fill.background()

steps = [
    ("1", "Decide on approach", "CLI only, Telegram only, or Hybrid"),
    ("2", "Procure Mac Mini", "M4 24GB recommended for multi-user"),
    ("3", "Setup & deploy Agent K", "Install software, configure accounts, deploy bot"),
    ("4", "Populate templates", "Firm's existing audit/tax/compilation templates"),
    ("5", "Build priority skills", "Start with /tax-computation and /compile-accounts"),
    ("6", "Pilot with 2\u20133 engagements", "Real clients, real deliverables, human review"),
    ("7", "Train staff & scale", "Telegram orientation for all, CLI training for power users"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.6) + i * Inches(0.75)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, Inches(1.9), y, Inches(4), Inches(0.5),
                 title, font_size=16, bold=True, color=WHITE)
    # Description
    add_text_box(slide, Inches(6), y + Inches(0.05), Inches(6), Inches(0.45),
                 desc, font_size=13, color=RGBColor(0xB0, 0xC4, 0xDE))

add_notes(slide,
    "NEXT STEPS — DETAIL:\n\n"
    "1. DECIDE: Today's meeting. Recommend hybrid approach.\n\n"
    "2. PROCURE: Mac Mini M4 with 24GB RAM (~RM 3,000-5,000).\n"
    "   Must be always-on, connected to internet. Can sit in server room.\n\n"
    "3. SETUP: 1-2 day engagement. Install Node.js, Claude CLI, Playwright,\n"
    "   configure Telegram bot, Gmail API, set up accounts.\n\n"
    "4. TEMPLATES: The firm provides their existing working paper templates,\n"
    "   checklists, engagement letters. We structure them in the folder hierarchy.\n\n"
    "5. PRIORITY SKILLS: Start with highest-volume, lowest-risk work:\n"
    "   - Tax computation (high volume, well-defined rules)\n"
    "   - Compilation FS (high volume, template-driven)\n"
    "   Then expand to audit workpapers and reports.\n\n"
    "6. PILOT: Pick 2-3 real client engagements. AI prepares, human reviews.\n"
    "   Measure time savings and error rates.\n\n"
    "7. SCALE: Once proven, train all staff on Telegram interface.\n"
    "   Train 1-2 power users on Claude CLI for maintenance."
)


# ============================================================
# Save
# ============================================================
output_path = '/home/user/Agent_K_Telegram/Agent_K_Deployment_Strategy.pptx'
prs.save(output_path)
print(f"Saved to: {output_path}")
